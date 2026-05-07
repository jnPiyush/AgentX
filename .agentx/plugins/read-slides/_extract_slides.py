"""AgentX read-slides extractor: pptx -> markdown."""

import sys, json, pathlib
from pptx import Presentation

src = sys.argv[1]
out = sys.argv[2]
prs = Presentation(src)
lines = ["# " + pathlib.Path(src).stem, ""]
for i, slide in enumerate(prs.slides, start=1):
    title = ""
    body = []
    notes = ""
    if slide.has_notes_slide:
        try:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes = ""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text:
            continue
        is_title = False
        try:
            is_title = shape == slide.shapes.title
        except Exception:
            pass
        if is_title and not title:
            title = text.splitlines()[0]
            extra = "\n".join(text.splitlines()[1:]).strip()
            if extra:
                body.append(extra)
        else:
            body.append(text)
    header = "## Slide {0}".format(i) + (": " + title if title else "")
    lines.append(header)
    lines.append("")
    for b in body:
        for ln in b.splitlines():
            ln = ln.rstrip()
            if not ln:
                continue
            if ln.startswith(("-", "*", "#")):
                lines.append(ln)
            else:
                lines.append("- " + ln)
        lines.append("")
    if notes:
        lines.append("> **Presenter notes**")
        for ln in notes.splitlines():
            ln = ln.strip()
            if ln:
                lines.append("> " + ln)
        lines.append("")
pathlib.Path(out).write_text("\n".join(lines), encoding="utf-8")
print(json.dumps({"slides": len(prs.slides), "out": out}))
