"""Build a visually rich AgentX pitch deck with native PowerPoint shapes."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Brand palette ---------------------------------------------------------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)  # background
DEEP = RGBColor(0x10, 0x2A, 0x4F)  # panel
CYAN = RGBColor(0x0E, 0xA5, 0xE9)  # accent (AgentX blue)
GREEN = RGBColor(0x22, 0xC5, 0x5E)  # accent (success)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)  # accent (warn)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)  # accent (skeptic)
INK = RGBColor(0xF8, 0xFA, 0xFC)  # primary text
MUTED = RGBColor(0x94, 0xA3, 0xB8)  # secondary text
LINE = RGBColor(0x33, 0x4E, 0x7A)  # divider

# 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.shadow.inherit = False
    return s


def text_box(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font="Segoe UI",
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def shape(slide, kind, x, y, w, h, fill=DEEP, line=LINE, line_w=0.75):
    sh = slide.shapes.add_shape(kind, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def header_bar(slide, eyebrow, title):
    # Cyan bar
    bar = shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0.6),
        Inches(0.55),
        Inches(0.12),
        Inches(0.55),
        fill=CYAN,
        line=CYAN,
    )
    text_box(
        slide,
        Inches(0.85),
        Inches(0.5),
        Inches(11),
        Inches(0.4),
        eyebrow.upper(),
        size=11,
        color=CYAN,
        bold=True,
    )
    text_box(
        slide,
        Inches(0.85),
        Inches(0.85),
        Inches(12),
        Inches(0.7),
        title,
        size=30,
        color=INK,
        bold=True,
    )
    # underline
    ul = shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0.85),
        Inches(1.55),
        Inches(11.6),
        Emu(9525),
        fill=LINE,
        line=LINE,
    )


def page_no(slide, n, total):
    text_box(
        slide,
        Inches(12.4),
        Inches(7.05),
        Inches(0.7),
        Inches(0.3),
        f"{n:02d} / {total:02d}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )
    text_box(
        slide,
        Inches(0.6),
        Inches(7.05),
        Inches(6),
        Inches(0.3),
        "AgentX  -  Digital Force for Software Delivery",
        size=10,
        color=MUTED,
    )


# ===========================================================================
# 1. TITLE
# ===========================================================================
def slide_title():
    s = add_slide()
    # accent stripe left
    shape(s, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH, fill=CYAN, line=CYAN)
    # Big logo dot
    dot = shape(
        s,
        MSO_SHAPE.OVAL,
        Inches(1.2),
        Inches(2.4),
        Inches(1.4),
        Inches(1.4),
        fill=CYAN,
        line=CYAN,
    )
    text_box(
        s,
        Inches(1.2),
        Inches(2.4),
        Inches(1.4),
        Inches(1.4),
        "X",
        size=72,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Title
    text_box(
        s,
        Inches(2.95),
        Inches(2.45),
        Inches(9),
        Inches(1.2),
        "AgentX",
        size=72,
        color=INK,
        bold=True,
    )
    text_box(
        s,
        Inches(2.95),
        Inches(3.55),
        Inches(9),
        Inches(0.5),
        "Digital Force for Software Delivery",
        size=22,
        color=CYAN,
    )
    text_box(
        s,
        Inches(2.95),
        Inches(4.15),
        Inches(10),
        Inches(1.4),
        "Turn AI coding agents into a structured engineering team\nwith routing, skills, memory, and multi-model deliberation.",
        size=16,
        color=MUTED,
    )
    # Footer
    text_box(
        s,
        Inches(1.2),
        Inches(6.6),
        Inches(11),
        Inches(0.4),
        "v8.4.47  -  Apache 2.0  -  github.com/jnPiyush/AgentX",
        size=12,
        color=MUTED,
    )


# ===========================================================================
# 2. THE PROBLEM
# ===========================================================================
def slide_problem(n, total):
    s = add_slide()
    header_bar(
        s, "The Problem", "Zero-shot AI coding does not scale to real engineering"
    )
    items = [
        ("87%", "of AI-generated code\nneeds rework before merge", AMBER),
        ("1 model", "single point of reasoning,\nsingle blind spot", ROSE),
        ("0 memory", "no decisions retained\nacross sessions", CYAN),
        ("No gate", "ships unreviewed,\nuntested, untraceable", GREEN),
    ]
    x0, y0, w, h, gap = (
        Inches(0.85),
        Inches(2.1),
        Inches(2.85),
        Inches(2.6),
        Inches(0.25),
    )
    for i, (big, small, col) in enumerate(items):
        x = x0 + (w + gap) * i
        card = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, w, h, fill=DEEP, line=LINE)
        text_box(
            s,
            x,
            y0 + Inches(0.35),
            w,
            Inches(1.0),
            big,
            size=44,
            color=col,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text_box(
            s,
            x + Inches(0.2),
            y0 + Inches(1.45),
            w - Inches(0.4),
            Inches(1.0),
            small,
            size=14,
            color=INK,
            align=PP_ALIGN.CENTER,
        )
    text_box(
        s,
        Inches(0.85),
        Inches(5.0),
        Inches(11.6),
        Inches(1.0),
        "Teams ship AI slop, then spend weeks cleaning it up.",
        size=20,
        color=MUTED,
        bold=False,
    )
    page_no(s, n, total)


# ===========================================================================
# 3. THE ANSWER
# ===========================================================================
def slide_answer(n, total):
    s = add_slide()
    header_bar(s, "The Answer", "AgentX is a harness, not a model")
    # Center concept
    text_box(
        s,
        Inches(0.85),
        Inches(2.0),
        Inches(11.6),
        Inches(0.6),
        "Force AI to plan, execute, iterate, review, and validate -- like a real team.",
        size=18,
        color=MUTED,
    )
    # 4 pillars
    pillars = [
        (
            "21",
            "Specialist Agents",
            "PM, Architect, UX,\nEngineer, Reviewer,\nTester, DevOps...",
            CYAN,
        ),
        (
            "94",
            "Production Skills",
            "Retrieval-led patterns\nloaded on demand,\nnot from model memory",
            GREEN,
        ),
        (
            "3+",
            "Model Council",
            "Diverse models debate\nhigh-stakes decisions\nbefore they ship",
            AMBER,
        ),
        (
            "\u221e",
            "Compound Memory",
            "Decisions, learnings,\nconventions persist\nin the repo",
            ROSE,
        ),
    ]
    x0, y, w, h, gap = (
        Inches(0.85),
        Inches(3.0),
        Inches(2.85),
        Inches(3.4),
        Inches(0.25),
    )
    for i, (big, label, body, col) in enumerate(pillars):
        x = x0 + (w + gap) * i
        card = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=DEEP, line=LINE)
        # accent strip top
        shape(s, MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.12), fill=col, line=col)
        text_box(
            s,
            x,
            y + Inches(0.4),
            w,
            Inches(0.9),
            big,
            size=48,
            color=col,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text_box(
            s,
            x,
            y + Inches(1.4),
            w,
            Inches(0.4),
            label,
            size=14,
            color=INK,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text_box(
            s,
            x + Inches(0.25),
            y + Inches(1.95),
            w - Inches(0.5),
            Inches(1.3),
            body,
            size=12,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
    page_no(s, n, total)


# ===========================================================================
# 4. THE TEAM (org chart diagram)
# ===========================================================================
def slide_team(n, total):
    s = add_slide()
    header_bar(
        s, "The AI Development Team", "21 specialists routed by Agent X (the hub)"
    )

    # Hub
    hub_w, hub_h = Inches(2.4), Inches(0.9)
    hx = (SW - hub_w) / 2
    hy = Inches(2.1)
    hub = shape(
        s, MSO_SHAPE.ROUNDED_RECTANGLE, hx, hy, hub_w, hub_h, fill=CYAN, line=CYAN
    )
    text_box(
        s,
        hx,
        hy,
        hub_w,
        hub_h,
        "Agent X (Hub)",
        size=18,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Branches
    branches = [
        ("Product & Design", "PM  -  UX Designer", GREEN),
        ("Architecture", "Architect  -  Data Scientist", CYAN),
        ("Engineering", "Engineer  -  DevOps", AMBER),
        ("Quality", "Reviewer  -  Tester  -  Auto-Fix", ROSE),
        ("Analytics & Gov.", "Power BI  -  Research  -  Coach", GREEN),
    ]
    bw, bh, gap = Inches(2.35), Inches(1.7), Inches(0.2)
    total_w = bw * 5 + gap * 4
    bx0 = (SW - total_w) / 2
    by = Inches(4.1)
    for i, (label, members, col) in enumerate(branches):
        bx = bx0 + i * (bw + gap)
        # connector line from hub bottom to branch top center
        line = shape(
            s,
            MSO_SHAPE.RECTANGLE,
            bx + bw / 2 - Emu(4500),
            hy + hub_h,
            Emu(9000),
            by - (hy + hub_h),
            fill=LINE,
            line=LINE,
        )
        # branch card
        card = shape(
            s, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh, fill=DEEP, line=LINE
        )
        shape(s, MSO_SHAPE.RECTANGLE, bx, by, bw, Inches(0.1), fill=col, line=col)
        text_box(
            s,
            bx,
            by + Inches(0.25),
            bw,
            Inches(0.4),
            label,
            size=13,
            color=col,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text_box(
            s,
            bx + Inches(0.15),
            by + Inches(0.75),
            bw - Inches(0.3),
            Inches(0.9),
            members,
            size=12,
            color=INK,
            align=PP_ALIGN.CENTER,
        )

    text_box(
        s,
        Inches(0.85),
        Inches(6.2),
        Inches(11.6),
        Inches(0.5),
        "Each agent has a strict role contract: deliverable, gates, self-review.",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    page_no(s, n, total)


# ===========================================================================
# 5. AGENTIC LOOP (flow diagram)
# ===========================================================================
def slide_loop(n, total):
    s = add_slide()
    header_bar(s, "The Agentic Loop", "Generate -> Verify -> Self-Review -> Done")

    steps = [
        ("Generate", "Write code\nagainst criteria", CYAN),
        ("Verify", "Run tests\nand lint", GREEN),
        ("Self-Review", "Structured findings:\nHIGH/MED/LOW", AMBER),
        ("Done", "All gates green,\nloop complete", ROSE),
    ]
    bw, bh = Inches(2.5), Inches(1.5)
    gap_x = Inches(0.3)
    arrow_w = Inches(0.55)
    total_w = bw * 4 + arrow_w * 3
    x = (SW - total_w) / 2
    y = Inches(2.6)

    for i, (title, body, col) in enumerate(steps):
        card = shape(
            s,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            bw,
            bh,
            fill=DEEP,
            line=col,
            line_w=1.5,
        )
        text_box(
            s,
            x,
            y + Inches(0.2),
            bw,
            Inches(0.5),
            title,
            size=18,
            color=col,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text_box(
            s,
            x,
            y + Inches(0.75),
            bw,
            Inches(0.7),
            body,
            size=12,
            color=INK,
            align=PP_ALIGN.CENTER,
        )
        if i < 3:
            ax = x + bw + Emu(50000)
            ay = y + bh / 2 - Inches(0.2)
            arr = shape(
                s,
                MSO_SHAPE.RIGHT_ARROW,
                ax,
                ay,
                arrow_w - Emu(100000),
                Inches(0.4),
                fill=CYAN,
                line=CYAN,
            )
            x = x + bw + arrow_w
        else:
            x = x + bw

    # Loop-back arrow under the row
    loop_y = y + bh + Inches(0.6)
    text_box(
        s,
        Inches(0.85),
        loop_y,
        Inches(11.6),
        Inches(0.4),
        "if findings -> Fix and Re-run  (min 5 iterations enforced by CLI)",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        bold=True,
    )

    # Stats row
    y2 = Inches(5.6)
    stats = [
        ("80%+", "coverage gate"),
        ("5", "min iterations"),
        ("0", "blocked secrets"),
        ("100%", "evidence-backed"),
    ]
    sw_, sh_, sg = Inches(2.7), Inches(1.0), Inches(0.25)
    sx = (SW - (sw_ * 4 + sg * 3)) / 2
    for i, (big, label) in enumerate(stats):
        x_ = sx + i * (sw_ + sg)
        c = shape(
            s, MSO_SHAPE.ROUNDED_RECTANGLE, x_, y2, sw_, sh_, fill=DEEP, line=LINE
        )
        text_box(
            s,
            x_,
            y2 + Inches(0.05),
            sw_,
            Inches(0.5),
            big,
            size=24,
            color=GREEN,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text_box(
            s,
            x_,
            y2 + Inches(0.55),
            sw_,
            Inches(0.4),
            label,
            size=12,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
    page_no(s, n, total)


# ===========================================================================
# 6. MODEL COUNCIL (triad diagram)
# ===========================================================================
def slide_council(n, total):
    s = add_slide()
    header_bar(s, "Model Council", "Stress-test the decision, not just the code")

    cx, cy = SW / 2, Inches(4.4)
    r = Inches(2.0)

    # Decision in center
    dw, dh = Inches(2.6), Inches(0.9)
    deci = shape(
        s,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        cx - dw / 2,
        cy - dh / 2,
        dw,
        dh,
        fill=NAVY,
        line=CYAN,
        line_w=2.0,
    )
    text_box(
        s,
        cx - dw / 2,
        cy - dh / 2,
        dw,
        dh,
        "Decision",
        size=18,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Three roles around the triangle
    roles = [
        ("Analyst", "Decompose\nevidence + facts", GREEN, (-1, -1)),
        ("Strategist", "Frame\nsecond-order effects", AMBER, (1, -1)),
        ("Skeptic", "Hunt\nfailure modes", ROSE, (0, 1)),
    ]
    rw, rh = Inches(2.6), Inches(1.4)
    for label, body, col, (dx, dy) in roles:
        rx = cx + dx * Inches(3.0) - rw / 2
        ry = cy + dy * Inches(2.2) - rh / 2
        card = shape(
            s,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            rx,
            ry,
            rw,
            rh,
            fill=DEEP,
            line=col,
            line_w=1.5,
        )
        text_box(
            s,
            rx,
            ry + Inches(0.2),
            rw,
            Inches(0.5),
            label,
            size=18,
            color=col,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text_box(
            s,
            rx,
            ry + Inches(0.75),
            rw,
            Inches(0.6),
            body,
            size=12,
            color=INK,
            align=PP_ALIGN.CENTER,
        )
        # connector to center
        line = shape(
            s,
            MSO_SHAPE.RECTANGLE,
            min(rx + rw / 2, cx) - Emu(4500),
            min(ry + rh / 2, cy),
            Emu(9000),
            abs((ry + rh / 2) - cy) or Emu(1),
            fill=LINE,
            line=LINE,
        )

    text_box(
        s,
        Inches(0.85),
        Inches(6.7),
        Inches(11.6),
        Inches(0.5),
        "Triggered for: PRD scope  -  ADR options  -  AI design  -  code review  -  research",
        size=13,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    page_no(s, n, total)


# ===========================================================================
# 7. WORKFLOW CHECKPOINTS (timeline)
# ===========================================================================
def slide_workflow(n, total):
    s = add_slide()
    header_bar(
        s, "Workflow Checkpoints", "Resolved from durable evidence, not chat history"
    )

    steps = [
        ("Brainstorm", "Frame", CYAN),
        ("Plan", "Anchor", GREEN),
        ("Work", "Build", AMBER),
        ("Review", "Settle", ROSE),
        ("Capture", "Compound", CYAN),
        ("Done", "Ship", GREEN),
    ]
    n_steps = len(steps)
    margin = Inches(0.9)
    track_w = SW - margin * 2
    y = Inches(3.6)
    # rail
    rail = shape(
        s,
        MSO_SHAPE.RECTANGLE,
        margin,
        y - Emu(4500),
        track_w,
        Emu(9000),
        fill=LINE,
        line=LINE,
    )
    spacing = track_w / (n_steps - 1)
    for i, (label, sub, col) in enumerate(steps):
        cx = margin + spacing * i
        # node
        d = Inches(0.55)
        node = shape(s, MSO_SHAPE.OVAL, cx - d / 2, y - d / 2, d, d, fill=col, line=col)
        text_box(
            s,
            cx - d / 2,
            y - d / 2,
            d,
            d,
            str(i + 1),
            size=16,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # label
        text_box(
            s,
            cx - Inches(1.0),
            y + Inches(0.5),
            Inches(2.0),
            Inches(0.4),
            label,
            size=15,
            color=INK,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text_box(
            s,
            cx - Inches(1.0),
            y + Inches(0.95),
            Inches(2.0),
            Inches(0.4),
            sub,
            size=11,
            color=col,
            align=PP_ALIGN.CENTER,
            font="Segoe UI",
        )
        # eyebrow above
        text_box(
            s,
            cx - Inches(1.0),
            y - Inches(1.0),
            Inches(2.0),
            Inches(0.4),
            f"Step {i+1}",
            size=10,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )

    text_box(
        s,
        Inches(0.85),
        Inches(5.8),
        Inches(11.6),
        Inches(0.5),
        "Evidence: execution plans, progress logs, review artifacts -- all repo-local",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    text_box(
        s,
        Inches(0.85),
        Inches(6.3),
        Inches(11.6),
        Inches(0.5),
        "Survives context loss. Resumable by any agent. Auditable forever.",
        size=14,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    page_no(s, n, total)


# ===========================================================================
# 8. SKILLS LIBRARY
# ===========================================================================
def slide_skills(n, total):
    s = add_slide()
    header_bar(
        s, "Skills Library", "Retrieval over recall  -  94 skills, 12 categories"
    )

    cats = [
        ("Architecture", "5", "api-design, security,\ndatabase, performance", CYAN),
        ("AI Systems", "29", "langgraph, foundry-sdk,\nrag, evaluation, safety", GREEN),
        ("Development", "15", "testing, error-handling,\niterative-loop, scrub", AMBER),
        ("Languages", "10", "C#, Python, TS, React,\nC, C++, Rust, Go", ROSE),
        ("Ops & Infra", "9", "GitHub Actions, Terraform,\nAzure, containers", CYAN),
        ("Data & Test", "13", "Databricks, Fabric,\ne2e, performance, security", GREEN),
        ("Design", "4", "design-system, ux-ui,\nprototype-craft, frontend", AMBER),
        ("Domain & PRD", "9", "finance, legal, oil&gas,\ntax, prd, diagrams", ROSE),
    ]
    cols, rows = 4, 2
    cw, ch = Inches(2.85), Inches(2.0)
    gx, gy = Inches(0.2), Inches(0.2)
    total_w = cw * cols + gx * (cols - 1)
    x0 = (SW - total_w) / 2
    y0 = Inches(2.1)
    for idx, (label, count, body, col) in enumerate(cats):
        r, c = divmod(idx, cols)
        x = x0 + c * (cw + gx)
        y = y0 + r * (ch + gy)
        card = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch, fill=DEEP, line=LINE)
        shape(s, MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), ch, fill=col, line=col)
        text_box(
            s,
            x + Inches(0.25),
            y + Inches(0.15),
            cw - Inches(0.4),
            Inches(0.5),
            label,
            size=15,
            color=INK,
            bold=True,
        )
        text_box(
            s,
            x + cw - Inches(0.85),
            y + Inches(0.1),
            Inches(0.7),
            Inches(0.5),
            count,
            size=22,
            color=col,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
        text_box(
            s,
            x + Inches(0.25),
            y + Inches(0.75),
            cw - Inches(0.4),
            Inches(1.2),
            body,
            size=11,
            color=MUTED,
        )
    text_box(
        s,
        Inches(0.85),
        Inches(6.65),
        Inches(11.6),
        Inches(0.5),
        "Max 3-4 skills per task (~20K tokens). Agents read, not guess.",
        size=14,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    page_no(s, n, total)


# ===========================================================================
# 9. PLUGINS
# ===========================================================================
def slide_plugins(n, total):
    s = add_slide()
    header_bar(s, "Plugins", "Markdown stays the source of truth")

    # Center markdown badge
    badge_w, badge_h = Inches(2.6), Inches(1.2)
    bx = (SW - badge_w) / 2
    by = Inches(3.5)
    badge = shape(
        s, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, badge_w, badge_h, fill=CYAN, line=CYAN
    )
    text_box(
        s,
        bx,
        by,
        badge_w,
        badge_h,
        ".md",
        size=44,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Outbound (right)
    out_items = [("convert-docs", ".docx"), ("convert-slides", ".pptx")]
    for i, (name, ext) in enumerate(out_items):
        bw, bh = Inches(2.5), Inches(1.0)
        x = Inches(9.6)
        y = Inches(2.6) + i * Inches(1.4)
        card = shape(
            s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, bh, fill=DEEP, line=GREEN
        )
        text_box(
            s,
            x + Inches(0.2),
            y + Inches(0.1),
            bw - Inches(0.4),
            Inches(0.4),
            name,
            size=13,
            color=GREEN,
            bold=True,
        )
        text_box(
            s,
            x + Inches(0.2),
            y + Inches(0.5),
            bw - Inches(0.4),
            Inches(0.4),
            f"-> {ext}",
            size=14,
            color=INK,
        )
        # arrow from md to plugin
        ax = bx + badge_w
        ay = by + badge_h / 2 - Inches(0.1)
        arrow = shape(
            s,
            MSO_SHAPE.RIGHT_ARROW,
            ax + Inches(0.1),
            y + bh / 2 - Inches(0.1),
            x - (ax + Inches(0.2)),
            Inches(0.2),
            fill=GREEN,
            line=GREEN,
        )

    # Inbound (left)
    in_items = [
        ("read-docs", "Word/RTF/HTML"),
        ("read-slides", "PowerPoint"),
        ("read-pdf", "PDF + page anchors"),
    ]
    for i, (name, ext) in enumerate(in_items):
        bw, bh = Inches(2.5), Inches(0.85)
        x = Inches(1.2)
        y = Inches(2.4) + i * Inches(1.1)
        card = shape(
            s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, bh, fill=DEEP, line=AMBER
        )
        text_box(
            s,
            x + Inches(0.2),
            y + Inches(0.05),
            bw - Inches(0.4),
            Inches(0.4),
            name,
            size=13,
            color=AMBER,
            bold=True,
        )
        text_box(
            s,
            x + Inches(0.2),
            y + Inches(0.45),
            bw - Inches(0.4),
            Inches(0.4),
            f"{ext} ->",
            size=12,
            color=INK,
        )
        # arrow from plugin to md
        ax = x + bw
        arrow = shape(
            s,
            MSO_SHAPE.RIGHT_ARROW,
            ax + Inches(0.1),
            y + bh / 2 - Inches(0.1),
            bx - (ax + Inches(0.2)),
            Inches(0.2),
            fill=AMBER,
            line=AMBER,
        )

    text_box(
        s,
        Inches(0.85),
        Inches(6.6),
        Inches(11.6),
        Inches(0.5),
        "Zero-install. Plugins shell out to Pandoc / Python on PATH only when used.",
        size=13,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    page_no(s, n, total)


# ===========================================================================
# 10. WHY IT WINS
# ===========================================================================
def slide_why(n, total):
    s = add_slide()
    header_bar(s, "Why AgentX Wins", "What changes when you adopt the harness")

    rows = [
        ("Without AgentX", "With AgentX", ROSE, GREEN),
        (
            "Ship AI slop, clean it up later",
            "Reviewed, tested, traceable on day one",
            None,
            None,
        ),
        (
            "Single-model blind spots",
            "Three-model Council debates each call",
            None,
            None,
        ),
        (
            "Decisions evaporate after the chat",
            "Plans, learnings, conventions in repo",
            None,
            None,
        ),
        (
            "Each engineer prompts differently",
            "21 agents with strict role contracts",
            None,
            None,
        ),
        (
            "Quality depends on the prompter",
            "Loop + 80% coverage + 5-iter gate",
            None,
            None,
        ),
    ]
    col_w = Inches(5.7)
    x_l = Inches(0.85)
    x_r = x_l + col_w + Inches(0.3)
    y0 = Inches(2.1)
    rh = Inches(0.7)
    for i, (left, right, lc, rc) in enumerate(rows):
        y = y0 + i * rh
        if i == 0:
            shape(
                s,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_l,
                y,
                col_w,
                rh - Inches(0.05),
                fill=ROSE,
                line=ROSE,
            )
            shape(
                s,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_r,
                y,
                col_w,
                rh - Inches(0.05),
                fill=GREEN,
                line=GREEN,
            )
            text_box(
                s,
                x_l,
                y,
                col_w,
                rh - Inches(0.05),
                left,
                size=16,
                color=NAVY,
                bold=True,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            text_box(
                s,
                x_r,
                y,
                col_w,
                rh - Inches(0.05),
                right,
                size=16,
                color=NAVY,
                bold=True,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
        else:
            shape(
                s,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_l,
                y,
                col_w,
                rh - Inches(0.05),
                fill=DEEP,
                line=LINE,
            )
            shape(
                s,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_r,
                y,
                col_w,
                rh - Inches(0.05),
                fill=DEEP,
                line=LINE,
            )
            # icons
            ic1 = shape(
                s,
                MSO_SHAPE.OVAL,
                x_l + Inches(0.2),
                y + Inches(0.18),
                Inches(0.3),
                Inches(0.3),
                fill=ROSE,
                line=ROSE,
            )
            text_box(
                s,
                x_l + Inches(0.2),
                y + Inches(0.18),
                Inches(0.3),
                Inches(0.3),
                "x",
                size=14,
                color=NAVY,
                bold=True,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            ic2 = shape(
                s,
                MSO_SHAPE.OVAL,
                x_r + Inches(0.2),
                y + Inches(0.18),
                Inches(0.3),
                Inches(0.3),
                fill=GREEN,
                line=GREEN,
            )
            text_box(
                s,
                x_r + Inches(0.2),
                y + Inches(0.18),
                Inches(0.3),
                Inches(0.3),
                "v",
                size=14,
                color=NAVY,
                bold=True,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            text_box(
                s,
                x_l + Inches(0.65),
                y,
                col_w - Inches(0.7),
                rh - Inches(0.05),
                left,
                size=14,
                color=INK,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            text_box(
                s,
                x_r + Inches(0.65),
                y,
                col_w - Inches(0.7),
                rh - Inches(0.05),
                right,
                size=14,
                color=INK,
                anchor=MSO_ANCHOR.MIDDLE,
            )
    page_no(s, n, total)


# ===========================================================================
# 11. CALL TO ACTION
# ===========================================================================
def slide_cta(n, total):
    s = add_slide()
    # full bleed accent
    shape(s, MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.5), fill=CYAN, line=CYAN)
    shape(
        s,
        MSO_SHAPE.RECTANGLE,
        0,
        SH - Inches(0.5),
        SW,
        Inches(0.5),
        fill=CYAN,
        line=CYAN,
    )

    text_box(
        s,
        Inches(0.85),
        Inches(1.6),
        Inches(11.6),
        Inches(0.6),
        "GET STARTED",
        size=14,
        color=CYAN,
        bold=True,
    )
    text_box(
        s,
        Inches(0.85),
        Inches(2.0),
        Inches(11.6),
        Inches(1.2),
        "Five minutes to your first reviewed feature.",
        size=36,
        color=INK,
        bold=True,
    )

    # install card
    code_box = shape(
        s,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.85),
        Inches(3.6),
        Inches(11.6),
        Inches(1.3),
        fill=DEEP,
        line=CYAN,
        line_w=1.5,
    )
    text_box(
        s,
        Inches(1.1),
        Inches(3.7),
        Inches(11),
        Inches(0.4),
        "PowerShell",
        size=10,
        color=MUTED,
        bold=True,
    )
    text_box(
        s,
        Inches(1.1),
        Inches(4.0),
        Inches(11.2),
        Inches(0.7),
        "irm https://raw.githubusercontent.com/jnPiyush/AgentX/v8.4.47/install.ps1 | iex",
        size=18,
        color=GREEN,
        bold=True,
        font="Cascadia Mono",
    )

    text_box(
        s,
        Inches(0.85),
        Inches(5.2),
        Inches(11.6),
        Inches(0.5),
        "or install the AgentX VS Code extension and run:",
        size=14,
        color=MUTED,
    )
    text_box(
        s,
        Inches(0.85),
        Inches(5.6),
        Inches(11.6),
        Inches(0.6),
        "AgentX: Initialize Local Runtime",
        size=22,
        color=CYAN,
        bold=True,
        font="Cascadia Mono",
    )

    text_box(
        s,
        Inches(0.85),
        Inches(6.4),
        Inches(11.6),
        Inches(0.5),
        "github.com/jnPiyush/AgentX  -  Apache 2.0  -  Open source",
        size=13,
        color=MUTED,
    )


# ===========================================================================
# Build deck
# ===========================================================================
slides = [
    slide_title,
    slide_problem,
    slide_answer,
    slide_team,
    slide_loop,
    slide_council,
    slide_workflow,
    slide_skills,
    slide_plugins,
    slide_why,
    slide_cta,
]
TOTAL = len(slides)
for i, fn in enumerate(slides, start=1):
    if i == 1:
        fn()
    else:
        fn(i, TOTAL)

import os

out = os.path.join("docs", "pitch", "agentx-pitch.pptx")
prs.save(out)
print(f"[OK] {out}  ({TOTAL} slides)")
